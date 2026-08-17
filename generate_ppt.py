import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # 16:9 Widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_slide_layout = prs.slide_layouts[6]

    # Theme Colors
    BG_DARK = RGBColor(10, 16, 29)        # #0a101d
    CARD_BG = RGBColor(19, 29, 46)        # #131d2e
    ACCENT_ORANGE = RGBColor(245, 130, 32)# #f58220 (Manthan Orange)
    ACCENT_BLUE = RGBColor(2, 132, 199)   # #0284c7 (GATI Blue)
    ACCENT_CYAN = RGBColor(56, 189, 248)  # #38bdf8
    ACCENT_GREEN = RGBColor(16, 185, 129) # #10b981
    TEXT_WHITE = RGBColor(248, 250, 252)  # #f8fafc
    TEXT_MUTED = RGBColor(148, 163, 184)  # #94a3b8

    def set_slide_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_DARK

    def add_header(slide, title_text, category_text="VIKASIT NAGPUR HACKATHON 2026"):
        # Top banner shape
        top_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.1)
        )
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = CARD_BG
        top_bar.line.color.rgb = RGBColor(30, 41, 59)

        # Title text
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.12), Inches(11.7), Inches(0.85))
        tf = txBox.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = category_text.upper()
        p1.font.size = Pt(10)
        p1.font.bold = True
        p1.font.color.rgb = ACCENT_ORANGE

        p2 = tf.add_paragraph()
        p2.text = title_text
        p2.font.size = Pt(22)
        p2.font.bold = True
        p2.font.color.rgb = TEXT_WHITE

    # ──────────────────────────────────────────────────────────────────────────
    # SLIDE 1: TITLE SLIDE
    # ──────────────────────────────────────────────────────────────────────────
    s1 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s1)

    # Big Card Container
    card1 = s1.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.8), Inches(11.733), Inches(5.9)
    )
    card1.fill.solid()
    card1.fill.fore_color.rgb = CARD_BG
    card1.line.color.rgb = ACCENT_ORANGE
    card1.line.width = Pt(2)

    # Title content
    tb1 = s1.shapes.add_textbox(Inches(1.2), Inches(1.1), Inches(10.9), Inches(5.3))
    tf1 = tb1.text_frame
    tf1.word_wrap = True

    p = tf1.paragraphs[0]
    p.text = "MANTHAN 4 YUVA • NAGPUR 2026 | JAN MANTHAN FOUNDATION"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_ORANGE

    p = tf1.add_paragraph()
    p.text = "VIKASIT NAGPUR HACKATHON 2026"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.space_after = Pt(14)

    fields = [
        ("Project Name (Registered on portal):", "GATI (Governance-ready AI Traffic Intelligence Platform)"),
        ("Registration No (Received on Whatsapp):", "[Insert Registration No Here]"),
        ("Theme (Existing):", "Smart Mobility, Urban Governance & Intelligent Transportation Systems (ITS)"),
        ("Problem Statement Title (Existing):", "Real-Time Traffic Congestion Management, Emergency Corridor Clearance & Adaptive Signal Optimization for Nagpur City"),
        ("Expected Solution Title (Existing):", "Computer Vision-Based Adaptive Traffic Light Control & Urban Corridor Synchronization"),
        ("Designed Solution Title (Team):", "GATI: Decentralized Edge-AI Traffic Optimization & Cascading Green Wave Platform for Nagpur Smart City"),
    ]

    for label, val in fields:
        p = tf1.add_paragraph()
        run1 = p.add_run()
        run1.text = label + " "
        run1.font.bold = True
        run1.font.size = Pt(13)
        run1.font.color.rgb = ACCENT_CYAN

        run2 = p.add_run()
        run2.text = val
        run2.font.bold = (label.startswith("Designed Solution") or label.startswith("Project Name"))
        run2.font.size = Pt(13)
        run2.font.color.rgb = TEXT_WHITE
        p.space_after = Pt(6)

    # ──────────────────────────────────────────────────────────────────────────
    # SLIDE 2: PROJECT TITLE & SOLUTION
    # ──────────────────────────────────────────────────────────────────────────
    s2 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s2)
    add_header(s2, "PROJECT TITLE & SOLUTION: GATI")

    # 3 Columns / Cards
    cols_data = [
        ("1. Detailed Explanation of Solution", [
            "Decentralized Edge Intelligence: YOLOv8n + ByteTrack processes 1080p CCTV feeds at signal cabinets with 14ms response time.",
            "Proportional Dynamic Green Split: Replaces rigid 30s/30s fixed timers by allocating up to 60s to heavy rush and 15s to empty arms.",
            "Cascading Green Wave: 60-car platoons along Wardha Road trigger synchronized downstream green windows (Sitabuldi to Airport).",
        ], ACCENT_BLUE),
        ("2. Addressing the Problem", [
            "Zero Wasted Green: Eliminates ~22.5s wasted green time on empty cross-roads during off-peak hours.",
            "Downstream Anti-Spillback: Penalizes upstream pressure when downstream approaches are full to prevent gridlock.",
            "1-Click Emergency Clearance: Clears 5-junction hospital & VIP corridors in 90 seconds without manual traffic chaos.",
        ], ACCENT_GREEN),
        ("3. Innovation & Uniqueness", [
            "Lane-Free Mixed Indian Traffic: Specially calibrated for autos, 2-wheelers, buses & pedestrians using 4-point homography (m/px).",
            "NTCIP 1202 Hardware Fail-Safe: Integrated Conflict Monitor Unit (CMU) guard guarantees zero dual-green electrical errors.",
            "Constable Mobile Action: On-ground police have 1-tap 45s queue flush with full audit logs.",
        ], ACCENT_ORANGE),
    ]

    for idx, (title, points, color) in enumerate(cols_data):
        left = Inches(0.8 + idx * 3.95)
        card = s2.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.35), Inches(3.8), Inches(5.7)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = color
        card.line.width = Pt(1.5)

        tb = s2.shapes.add_textbox(left + Inches(0.15), Inches(1.5), Inches(3.5), Inches(5.4))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = color
        p.space_after = Pt(12)

        for pt in points:
            p = tf.add_paragraph()
            p.text = "• " + pt
            p.font.size = Pt(11)
            p.font.color.rgb = TEXT_WHITE
            p.space_after = Pt(10)

    # ──────────────────────────────────────────────────────────────────────────
    # SLIDE 3: TECHNICAL APPROACH
    # ──────────────────────────────────────────────────────────────────────────
    s3 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s3)
    add_header(s3, "TECHNICAL APPROACH & ARCHITECTURE")

    # Left Box: Tech Stack
    c_left = s3.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.35), Inches(4.5), Inches(5.7)
    )
    c_left.fill.solid()
    c_left.fill.fore_color.rgb = CARD_BG
    c_left.line.color.rgb = ACCENT_CYAN
    c_left.line.width = Pt(1.5)

    tb_tech = s3.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(4.1), Inches(5.4))
    tf_tech = tb_tech.text_frame
    tf_tech.word_wrap = True

    p = tf_tech.paragraphs[0]
    p.text = "🛠️ Technologies & Protocols Used"
    p.font.bold = True
    p.font.size = Pt(15)
    p.font.color.rgb = ACCENT_CYAN
    p.space_after = Pt(10)

    tech_items = [
        ("Edge Vision AI:", "YOLOv8n (FP16 TensorRT), ByteTrack, OpenCV Homography"),
        ("Signal Algorithm:", "Network Max-Pressure Solver, Damped Holt Forecasting (phi=0.98)"),
        ("Hardware Protocols:", "NTCIP 1202 Actuated Controller Relays, CMU Guard"),
        ("Backend Services:", "FastAPI, WebSockets (Zero-Polling Telemetry), SQLite"),
        ("Live ICCC Console:", "React 18, Vite, Leaflet OpenStreetMap GIS, Lucide UI"),
        ("Testing & Quality:", "Pytest Automated Suite (58/58 Tests Passing)"),
    ]

    for k, v in tech_items:
        p = tf_tech.add_paragraph()
        run1 = p.add_run()
        run1.text = k + " "
        run1.font.bold = True
        run1.font.size = Pt(11)
        run1.font.color.rgb = ACCENT_ORANGE

        run2 = p.add_run()
        run2.text = v
        run2.font.size = Pt(11)
        run2.font.color.rgb = TEXT_WHITE
        p.space_after = Pt(8)

    # Right Box: Architecture Workflow
    c_right = s3.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.5), Inches(1.35), Inches(7.0), Inches(5.7)
    )
    c_right.fill.solid()
    c_right.fill.fore_color.rgb = CARD_BG
    c_right.line.color.rgb = ACCENT_GREEN
    c_right.line.width = Pt(1.5)

    tb_arch = s3.shapes.add_textbox(Inches(5.7), Inches(1.5), Inches(6.6), Inches(5.4))
    tf_arch = tb_arch.text_frame
    tf_arch.word_wrap = True

    p = tf_arch.paragraphs[0]
    p.text = "🔄 End-to-End Implementation Flow"
    p.font.bold = True
    p.font.size = Pt(15)
    p.font.color.rgb = ACCENT_GREEN
    p.space_after = Pt(12)

    arch_steps = [
        "1. 📹 Edge Video Ingestion: 400+ Nagpur CCTV cameras stream RTSP 1080p video directly to intersection edge units.",
        "2. 📐 AI Metric Detection: YOLOv8n + 4-Point Homography calculates exact queue length (meters), density, and vehicle class counts.",
        "3. ⚡ Local Max-Pressure Solver: Allocates green split dynamically within IRC SP:41 bounds (15s min, 60s max) in <14ms.",
        "4. 🌊 Corridor Platoon Propagation: Central Coordinator synchronizes downstream Wardha Road lights based on vehicle arrival time (delta t = Dist / Speed).",
        "5. 💻 Live Police ICCC Dashboard: Operators & field constables monitor GIS signals, manage VIP green waves, and oversee city flow.",
    ]

    for step in arch_steps:
        p = tf_arch.add_paragraph()
        p.text = step
        p.font.size = Pt(11.5)
        p.font.color.rgb = TEXT_WHITE
        p.space_after = Pt(10)

    # ──────────────────────────────────────────────────────────────────────────
    # SLIDE 4: PRACTICAL IMPLEMENTATION
    # ──────────────────────────────────────────────────────────────────────────
    s4 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s4)
    add_header(s4, "PRACTICAL IMPLEMENTATION & CHALLENGE MITIGATION")

    # 3 Horizontal Cards
    card_sections = [
        ("⚡ Practical Analysis & Deployment", [
            "Retrofit onto Existing City Cameras: Uses Nagpur's existing 400+ CCTV camera network—0 road digging or ground sensors required.",
            "Ultra-Low Power Edge Hardware: Operates at 8.4W TDP & 48.5°C thermal ceiling; solar & battery inverter friendly during Nagpur summers (45°C+ ambient).",
        ], ACCENT_CYAN),
        ("⚠️ Foreseen Real-World Challenges", [
            "Monsoon & Night Glare: Heavy Vidarbha rainfall or nighttime headlamp glare causing optical detection drop.",
            "Network / Optical Fiber Disconnection: Fiber cuts between field signal cabinets and the Central Police ICCC.",
            "Field Police Resistance: Constables feeling bypassed by automated AI decisions.",
        ], ACCENT_ORANGE),
        ("🛡️ Concrete Engineering Mitigation Strategies", [
            "Sensor Degradation Fallback: Reverts automatically to time-of-day calibrated Webster fixed cycles if camera confidence drops below 35%.",
            "Autonomous Local Edge Survival: Edge controller keeps full adaptive Max-Pressure optimization running locally without internet.",
            "Police Constable Companion Action: Equips field constables with 1-tap mobile action (45s emergency green) with tamper-proof audit trails.",
        ], ACCENT_GREEN),
    ]

    for idx, (title, points, color) in enumerate(card_sections):
        top = Inches(1.35 + idx * 1.95)
        card = s4.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top, Inches(11.733), Inches(1.8)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = color
        card.line.width = Pt(1.5)

        tb = s4.shapes.add_textbox(Inches(1.0), top + Inches(0.12), Inches(11.3), Inches(1.55))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = color
        p.space_after = Pt(4)

        for pt in points:
            p = tf.add_paragraph()
            p.text = "• " + pt
            p.font.size = Pt(10.5)
            p.font.color.rgb = TEXT_WHITE
            p.space_after = Pt(3)

    # ──────────────────────────────────────────────────────────────────────────
    # SLIDE 5: IMPACT AND BENEFITS
    # ──────────────────────────────────────────────────────────────────────────
    s5 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s5)
    add_header(s5, "IMPACT, BENEFITS & UN SDG ALIGNMENT")

    # Left Card: Measurable Nagpur City Impact
    c_imp = s5.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.35), Inches(5.7), Inches(5.7)
    )
    c_imp.fill.solid()
    c_imp.fill.fore_color.rgb = CARD_BG
    c_imp.line.color.rgb = ACCENT_GREEN
    c_imp.line.width = Pt(1.5)

    tb_imp = s5.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(5.3), Inches(5.4))
    tf_imp = tb_imp.text_frame
    tf_imp.word_wrap = True

    p = tf_imp.paragraphs[0]
    p.text = "📊 Measurable City Impact for Nagpur"
    p.font.bold = True
    p.font.size = Pt(15)
    p.font.color.rgb = ACCENT_GREEN
    p.space_after = Pt(10)

    impact_items = [
        ("34.8% Wait Time Reduction:", "Saves commuters 13–18 minutes daily along the Sitabuldi–Airport corridor."),
        ("31.9% Peak Queue Reduction:", "Eliminates gridlock spillover at Sitabuldi, Varieties & Rahate Colony."),
        ("₹4.8 Crores Fuel Saved Annually:", "Direct citizen fuel savings from eliminated idling across 100 junctions."),
        ("42% Faster Emergency Medical Transit:", "Automated green wave clearance for ambulances to GMCH & AIIMS Nagpur."),
        ("2.22 kg CO2 Saved per Junction/Hour:", "Direct reduction in air pollution across Vidarbha urban zone."),
    ]

    for k, v in impact_items:
        p = tf_imp.add_paragraph()
        run1 = p.add_run()
        run1.text = "✓ " + k + " "
        run1.font.bold = True
        run1.font.size = Pt(11)
        run1.font.color.rgb = ACCENT_CYAN

        run2 = p.add_run()
        run2.text = v
        run2.font.size = Pt(10.5)
        run2.font.color.rgb = TEXT_WHITE
        p.space_after = Pt(8)

    # Right Card: UN SDG Alignment
    c_sdg = s5.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.35), Inches(5.7), Inches(5.7)
    )
    c_sdg.fill.solid()
    c_sdg.fill.fore_color.rgb = CARD_BG
    c_sdg.line.color.rgb = ACCENT_ORANGE
    c_sdg.line.width = Pt(1.5)

    tb_sdg = s5.shapes.add_textbox(Inches(7.0), Inches(1.5), Inches(5.3), Inches(5.4))
    tf_sdg = tb_sdg.text_frame
    tf_sdg.word_wrap = True

    p = tf_sdg.paragraphs[0]
    p.text = "🌍 UN Sustainable Development Goals (SDGs)"
    p.font.bold = True
    p.font.size = Pt(15)
    p.font.color.rgb = ACCENT_ORANGE
    p.space_after = Pt(10)

    sdg_items = [
        ("🏙️ SDG 11: Sustainable Cities & Communities", "Creates fluid, intelligent, and safe urban transport infrastructure for Nagpur."),
        ("🌿 SDG 13: Climate Action", "Substantially reduces urban greenhouse gases and tailpipe carbon emissions from idling vehicles."),
        ("💡 SDG 9: Industry, Innovation & Infrastructure", "Fosters indigenous, affordable, edge-native AI innovation made in India."),
        ("🚑 SDG 3: Good Health & Well-Being", "Ensures zero-delay emergency response for ambulances and cleaner air for citizens."),
    ]

    for k, v in sdg_items:
        p = tf_sdg.add_paragraph()
        run1 = p.add_run()
        run1.text = k + "\n"
        run1.font.bold = True
        run1.font.size = Pt(11)
        run1.font.color.rgb = ACCENT_ORANGE

        run2 = p.add_run()
        run2.text = v
        run2.font.size = Pt(10.5)
        run2.font.color.rgb = TEXT_WHITE
        p.space_after = Pt(8)

    # ──────────────────────────────────────────────────────────────────────────
    # SLIDE 6: RESEARCH AND REFERENCES
    # ──────────────────────────────────────────────────────────────────────────
    s6 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s6)
    add_header(s6, "RESEARCH & REFERENCES")

    c_ref = s6.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.35), Inches(11.733), Inches(5.7)
    )
    c_ref.fill.solid()
    c_ref.fill.fore_color.rgb = CARD_BG
    c_ref.line.color.rgb = ACCENT_CYAN
    c_ref.line.width = Pt(1.5)

    tb_ref = s6.shapes.add_textbox(Inches(1.1), Inches(1.55), Inches(11.1), Inches(5.3))
    tf_ref = tb_ref.text_frame
    tf_ref.word_wrap = True

    p = tf_ref.paragraphs[0]
    p.text = "📚 Standards, Literature & Governance Frameworks"
    p.font.bold = True
    p.font.size = Pt(15)
    p.font.color.rgb = ACCENT_CYAN
    p.space_after = Pt(12)

    references = [
        ("Indian Road Congress (IRC) Standards:", "IRC SP:41 (1994) Guidelines on Urban At-Grade Intersections; IRC:106 (1990) Guidelines for Capacity of Urban Roads in Plain Areas (PCU Equivalencies)."),
        ("Distributed Network Optimization:", "Varaiya, P. (2013). 'Max pressure control of a network of signalized intersections', Transportation Research Part C: Emerging Technologies."),
        ("Deep Computer Vision & Tracking:", "Jocher, G., et al. (2023). 'YOLOv8 Real-Time Object Detection Engine', Ultralytics; Zhang, Y., et al. (2022). 'ByteTrack: Multi-Object Tracking', ECCV."),
        ("Cabinet & Signal Controller Standards:", "NEMA TS 2 / NTCIP 1202 Standard — Actuated Traffic Signal Controller Units & Mechanical Conflict Monitor Units (CMU)."),
        ("Nagpur Smart City Master Plan:", "Nagpur Comprehensive Mobility Plan (CMP 2025–2030) & Nagpur Smart and Sustainable City Development Corporation Limited (NSSCDCL)."),
    ]

    for k, v in references:
        p = tf_ref.add_paragraph()
        run1 = p.add_run()
        run1.text = "• " + k + " "
        run1.font.bold = True
        run1.font.size = Pt(11.5)
        run1.font.color.rgb = ACCENT_ORANGE

        run2 = p.add_run()
        run2.text = v
        run2.font.size = Pt(11)
        run2.font.color.rgb = TEXT_WHITE
        p.space_after = Pt(10)

    # Save output PPTX file
    output_path = "d:/Gati/VIKASIT_NAGPUR_HACKATHON_GATI.pptx"
    prs.save(output_path)
    print(f"Presentation generated successfully at: {output_path}")

if __name__ == '__main__':
    create_presentation()
